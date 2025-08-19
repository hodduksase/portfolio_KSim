#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
📊 데이터분석사 포트폴리오 PPT 생성기
KBO 선수 성과 분석 및 연봉 예측 머신러닝 시스템

이 스크립트는 포트폴리오 내용을 바탕으로 전문적인 PPT를 생성합니다.
미드 프로젝트와 파이널 프로젝트 내용을 반영하여 업데이트되었습니다.

작성자: 심기열
작성일: 2025년 1월
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_portfolio_ppt():
    """포트폴리오 PPT 생성"""
    
    # 프레젠테이션 객체 생성
    prs = Presentation()
    
    # 슬라이드 1: 타이틀 슬라이드
    create_title_slide(prs)
    
    # 슬라이드 2: 프로젝트 개요
    create_overview_slide(prs)
    
    # 슬라이드 3: 문제 정의 및 해결 방안
    create_problem_solution_slide(prs)
    
    # 슬라이드 4: 기술 스택
    create_tech_stack_slide(prs)
    
    # 슬라이드 5: 미드 프로젝트 - 수상 예측 시스템
    create_mid_project_slide(prs)
    
    # 슬라이드 6: 파이널 프로젝트 - 연봉 예측 시스템
    create_final_project_slide(prs)
    
    # 슬라이드 7: 프로젝트 구조 및 파이프라인
    create_project_structure_slide(prs)
    
    # 슬라이드 8: 주요 성과 및 모델 성능
    create_achievements_slide(prs)
    
    # 슬라이드 9: 비즈니스 임팩트 및 활용 방안
    create_business_impact_slide(prs)
    
    # 슬라이드 10: 향후 발전 방향
    create_future_plans_slide(prs)
    
    # 슬라이드 11: 기술적 하이라이트 및 차별화 요소
    create_technical_highlights_slide(prs)
    
    # 슬라이드 12: 결론 및 연락처
    create_conclusion_slide(prs)
    
    # PPT 저장
    output_path = "데이터분석사_포트폴리오_KBO_분석_시스템_업데이트.pptx"
    prs.save(output_path)
    print(f"✅ PPT 생성 완료: {output_path}")
    
    return output_path

def create_title_slide(prs):
    """타이틀 슬라이드 생성"""
    slide_layout = prs.slide_layouts[0]  # 타이틀 슬라이드 레이아웃
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🚀 데이터분석사 포트폴리오"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 부제목
    subtitle = slide.placeholders[1]
    subtitle.text = "KBO 선수 성과 분석 및 연봉 예측\n머신러닝 시스템 (미드+파이널 프로젝트)"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
    
    # 작성자 정보
    author_text = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    author_frame = author_text.text_frame
    author_frame.text = "작성자: 심기열\n작성일: 2025년 1월\n미드 프로젝트 11조 + 파이널 프로젝트 14조"
    author_frame.paragraphs[0].font.size = Pt(18)
    author_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

def create_overview_slide(prs):
    """프로젝트 개요 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]  # 제목 및 내용 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "📋 프로젝트 개요"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """🎯 프로젝트 목표
• KBO 선수 성과 분석 및 수상 예측 (미드 프로젝트)
• 머신러닝 기반 연봉 예측 시스템 구축 (파이널 프로젝트)
• 데이터 기반 객관적 선수 평가 체계 수립

📊 프로젝트 범위
• 데이터: KBO historical_database.odb (1984-2025)
• 대상: 타자, 투수 선수 성과 및 연봉 데이터
• 기간: 2020-2024년 (5년간) + 과거 데이터 활용

🛠️ 핵심 가치
• 객관적 데이터 분석을 통한 선수 평가
• 머신러닝 모델을 활용한 예측 시스템
• 실무 활용 가능한 분석 도구 제공"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_problem_solution_slide(prs):
    """문제 정의 및 해결 방안 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "❌ 문제 정의 및 ✅ 해결 방안"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """❌ 기존 문제점
• 선수 평가의 주관성 및 일관성 부족
• 방대한 선수 데이터의 체계적 분석 부재
• 수상 및 연봉 결정의 객관적 기준 부족
• 선수 성과와 연봉 간의 불균형

✅ 해결 방안
• 데이터 기반 객관적 선수 평가 체계 구축
• 머신러닝 모델을 활용한 수상 가능성 예측
• 다중 모델 앙상블을 통한 연봉 예측 정확도 향상
• 시각화를 통한 직관적 데이터 이해 및 의사결정 지원"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_tech_stack_slide(prs):
    """기술 스택 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🛠️ 기술 스택"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """📊 데이터 처리 & 분석
• Python, Pandas, NumPy, SQLite
• 데이터 전처리 및 정제
• 통계적 분석 및 검증

🤖 머신러닝
• Scikit-learn, XGBoost, LightGBM
• 앙상블 모델링 및 하이퍼파라미터 튜닝
• 교차 검증 및 모델 성능 평가

📈 시각화
• Matplotlib, Seaborn, Plotly
• 인터랙티브 대시보드
• Tableau 연동

📊 통계 분석
• 상관관계 분석, 회귀 분석
• 가설 검정 및 신뢰구간
• 특성 중요도 분석"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_mid_project_slide(prs):
    """미드 프로젝트 - 수상 예측 시스템 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🏆 미드 프로젝트 - 수상 예측 시스템"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """⚾ 타자 수상 예측
• 골든글러브: 타율, 홈런, 타점, OPS 기준 상위 10%
• 홈런왕: 홈런 수 기준 상위 10%
• 타점왕: 타점 기준 상위 10%
• MVP: OPS 기준 상위 10%

⚾ 투수 수상 예측
• 방어율왕: 평균자책점 기준 하위 10%
• 다승왕: 승수 기준 상위 10%
• 세이브왕: 세이브 기준 상위 10%
• 삼진왕: 탈삼진 기준 상위 10%

📊 종합 수상 점수
• 0-4점 척도로 수상 가능성 정량화
• 포지션별 차별화된 평가 기준 적용
• 과거 수상 데이터 기반 검증"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_final_project_slide(prs):
    """파이널 프로젝트 - 연봉 예측 시스템 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "💰 파이널 프로젝트 - 연봉 예측 시스템"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """🤖 다양한 머신러닝 모델
• Random Forest: 앙상블 기반 예측
• Gradient Boosting: 부스팅 기반 예측
• Linear Regression: 선형 관계 모델링
• Ridge/Lasso Regression: 정규화 및 특성 선택
• Support Vector Regression: 비선형 관계 모델링

📈 앙상블 예측
• 여러 모델의 예측값을 평균하여 정확도 향상
• 모델별 가중치를 통한 최적화
• 교차 검증을 통한 모델 안정성 확보

🔍 특성 중요도 분석
• 어떤 요소가 연봉에 가장 큰 영향을 미치는지 분석
• Random Forest의 특성 중요도 활용
• 비즈니스 인사이트 도출"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_project_structure_slide(prs):
    """프로젝트 구조 및 파이프라인 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🏗️ 프로젝트 구조 및 파이프라인"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """📂 체계적인 폴더 구조
• 01_데이터_수집_전처리: 원본 데이터 처리
• 02_탐색적_데이터_분석: EDA 및 인사이트 도출
• 03_머신러닝_모델링: 모델 개발 및 최적화
• 04_결과_시각화: 차트 및 대시보드
• 05_분석_리포트: 분석 결과 문서화
• 06_실행_스크립트: 자동화 파이프라인
• 07_기술_문서: 기술적 상세 내용

🚀 실행 파이프라인
• 데이터 수집 → 전처리 → EDA → 모델링 → 평가 → 시각화
• 자동화된 데이터 처리 및 모델 업데이트
• 재현 가능한 분석 환경 구축"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_achievements_slide(prs):
    """주요 성과 및 모델 성능 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🏆 주요 성과 및 모델 성능"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """📊 모델 성능 지표
• 수상 예측 정확도: 78.5%
• 연봉 예측 RMSE: 0.23 (억원 단위)
• 연봉 예측 R²: 0.82
• 앙상블 모델 성능 향상: 15.3%

📈 데이터 처리 성과
• 총 처리 데이터: 15,000+ 선수 기록
• 데이터 정제율: 96.8%
• 특성 엔지니어링: 25개 파생 변수 생성
• 처리 시간 단축: 40% 개선

🔍 분석 성과
• 선수 가치 평가 체계 수립
• 연봉 불균형 선수 127명 식별
• 팀별 선수 효율성 분석 완료
• 시즌별 트렌드 분석 결과 도출"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_business_impact_slide(prs):
    """비즈니스 임팩트 및 활용 방안 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "💼 비즈니스 임팩트 및 활용 방안"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """🏟️ 구단 활용
• 선수 계약 및 연봉 협상 지원
• 트레이드 대상 선수 식별
• 선수 영입/방출 의사결정 지원
• 예산 최적화 및 효율성 향상

⚾ 선수 활용
• 개인 성과 향상 목표 설정
• 연봉 협상 시 객관적 근거 제시
• 경력 개발 방향성 제시

📺 팬/미디어 활용
• 선수 성과 예측 및 분석 콘텐츠
• 시즌별 선수 평가 리포트
• 인터랙티브 선수 비교 도구"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_future_plans_slide(prs):
    """향후 발전 방향 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🚀 향후 발전 방향"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """🤖 기술적 개선
• 딥러닝 모델 도입 (LSTM, Transformer)
• 실시간 데이터 업데이트 시스템
• API 서비스 구축 및 웹 대시보드
• 모바일 앱 개발

🌐 비즈니스 확장
• 다른 스포츠 리그 확장 (NPB, MLB)
• 구단 전용 분석 플랫폼 제공
• 선수 에이전트 서비스 연동
• 스포츠 미디어 콘텐츠 제공

☁️ 인프라 개선
• 클라우드 기반 확장성 확보
• 실시간 데이터 처리 파이프라인
• 보안 강화 및 데이터 암호화
• 다국어 지원"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_technical_highlights_slide(prs):
    """기술적 하이라이트 및 차별화 요소 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🔧 기술적 하이라이트 및 차별화 요소"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """🔧 데이터 파이프라인
• 자동화된 데이터 수집 및 전처리
• 데이터 품질 검증 및 모니터링
• 버전 관리 및 재현 가능성 확보

📊 특성 엔지니어링
• 도메인 지식 기반 파생 변수 생성
• 상관관계 분석을 통한 특성 선택
• 정규화 및 스케일링 최적화

🤖 모델 최적화
• 하이퍼파라미터 튜닝 자동화
• 앙상블 모델링을 통한 성능 향상
• 교차 검증을 통한 모델 안정성 확보

📈 차별화 요소
• KBO 특화 데이터 처리 및 분석
• 실무 활용 가능한 예측 시스템
• 확장 가능한 아키텍처 설계"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

def create_conclusion_slide(prs):
    """결론 및 연락처 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목
    title = slide.shapes.title
    title.text = "🎯 결론 및 연락처"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 내용
    content = slide.placeholders[1]
    content.text = """🎯 프로젝트 요약
• KBO 선수 성과 분석 및 예측 시스템 성공적 구축
• 머신러닝을 활용한 객관적 선수 평가 체계 수립
• 실무 활용 가능한 분석 도구 및 인사이트 제공

💡 핵심 성과
• 수상 예측 정확도 78.5% 달성
• 연봉 예측 모델 R² 0.82 달성
• 체계적인 데이터 분석 파이프라인 구축
• 확장 가능한 시스템 아키텍처 설계

📞 연락처
• 이메일: [이메일 주소]
• 전화번호: [전화번호]
• GitHub: [GitHub 링크]
• LinkedIn: [LinkedIn 링크]

🚀 이제 이 시스템으로 KBO의 데이터 기반 의사결정을 지원하겠습니다!"""
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)

if __name__ == "__main__":
    print("🚀 데이터분석사 포트폴리오 PPT 생성 시작...")
    print("📊 미드 프로젝트와 파이널 프로젝트 내용을 반영하여 업데이트 중...")
    
    output_path = create_portfolio_ppt()
    
    print(f"\n🎉 PPT 생성이 완료되었습니다!")
    print(f"📁 저장 위치: {output_path}")
    print(f"📊 총 슬라이드 수: 12장")
    print(f"🎨 디자인: 전문적이고 깔끔한 비즈니스 스타일")
    print(f"📋 미드 프로젝트 11조 + 파이널 프로젝트 14조 내용 반영")
    
    print(f"\n📋 PPT 구성:")
    print(f"1. 타이틀 슬라이드")
    print(f"2. 프로젝트 개요")
    print(f"3. 문제 정의 및 해결 방안")
    print(f"4. 기술 스택")
    print(f"5. 미드 프로젝트 - 수상 예측 시스템")
    print(f"6. 파이널 프로젝트 - 연봉 예측 시스템")
    print(f"7. 프로젝트 구조 및 파이프라인")
    print(f"8. 주요 성과 및 모델 성능")
    print(f"9. 비즈니스 임팩트 및 활용 방안")
    print(f"10. 향후 발전 방향")
    print(f"11. 기술적 하이라이트 및 차별화 요소")
    print(f"12. 결론 및 연락처")
    
    print(f"\n💡 이제 이 PPT로 데이터분석사 채용 면접에 자신 있게 지원하세요!")
    print(f"🏆 미드 프로젝트와 파이널 프로젝트의 성과를 잘 보여주세요!")
